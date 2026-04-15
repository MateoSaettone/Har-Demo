"""Dump the structure of the source deck so we can match its style."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
import json

HERE = Path(__file__).resolve().parent
src = HERE / "source.pptx"
prs = Presentation(src)

def fmt_color(c):
    try:
        if c.type is None:
            return None
        # RGBColor
        return str(c.rgb) if c.rgb is not None else f"theme:{c.theme_color}"
    except Exception:
        return None

print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU "
      f"({Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in)")
print(f"Slide layouts: {len(prs.slide_layouts)}")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  layout[{i}]: {layout.name}")

print(f"\nSlides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} (layout: {slide.slide_layout.name}) ---")
    for sh in slide.shapes:
        kind = sh.shape_type
        name = sh.name
        try:
            txt = sh.text_frame.text if sh.has_text_frame else ""
        except Exception:
            txt = ""
        txt = txt.replace("\n", " | ")[:160]
        fill = ""
        try:
            if sh.fill.type is not None:
                fill = f" fill={fmt_color(sh.fill.fore_color)}"
        except Exception:
            pass
        print(f"  [{kind}] {name}{fill} :: {txt!r}")
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs[:3]:
                for r in p.runs[:3]:
                    font = r.font
                    color = None
                    try:
                        color = fmt_color(font.color)
                    except Exception:
                        pass
                    print(f"      run: size={font.size} bold={font.bold} "
                          f"name={font.name} color={color} text={r.text!r}")
