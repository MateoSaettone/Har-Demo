"""
Build the final presentation from source.pptx.

- Fill empty placeholders (Our System, Training Pipeline, Limitations, Conclusion).
- Insert two new Results slides after Training Pipeline (scores + confusion matrix).
- Emit CIS4930 UG Group 3 Project Presentation - FINAL.pptx into this folder
  and copy it into the user's Downloads directory.
"""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import shutil

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
REPO = HERE.parent
SRC = HERE / "source.pptx"
OUT = HERE / "CIS4930 UG Group 3 Project Presentation - FINAL.pptx"
DOWNLOADS_OUT = Path.home() / "Downloads" / "CIS4930 UG Group 3 Project Presentation - FINAL.pptx"

LABELS = [
    "WALKING",
    "WALK\nUPSTAIRS",
    "WALK\nDOWNSTAIRS",
    "SITTING",
    "STANDING",
    "LAYING",
]

CNN_CM = np.array([
    [496, 0, 0, 0, 0, 0],
    [10, 437, 22, 0, 2, 0],
    [12, 41, 367, 0, 0, 0],
    [0, 6, 0, 368, 117, 0],
    [1, 0, 0, 67, 464, 0],
    [0, 2, 0, 0, 0, 535],
])
RF_CM = np.array([
    [437, 19, 40, 0, 0, 0],
    [70, 376, 24, 1, 0, 0],
    [23, 20, 377, 0, 0, 0],
    [3, 23, 0, 385, 80, 0],
    [2, 4, 0, 123, 403, 0],
    [0, 0, 0, 0, 0, 537],
])

CNN_ACC = 0.9050
RF_ACC = 0.8534

CNN_PR = {
    "WALKING":            (0.9557, 1.0000, 0.9773),
    "WALKING_UPSTAIRS":   (0.8992, 0.9278, 0.9133),
    "WALKING_DOWNSTAIRS": (0.9434, 0.8738, 0.9073),
    "SITTING":            (0.8460, 0.7495, 0.7948),
    "STANDING":           (0.7959, 0.8722, 0.8323),
    "LAYING":             (1.0000, 0.9963, 0.9981),
}


# --------------------------------------------------------------------- images

def plot_confusion_matrix(cm, title, fname):
    cm = cm.astype(float)
    totals = cm.sum(axis=1, keepdims=True)
    totals[totals == 0] = 1
    norm = cm / totals
    fig, ax = plt.subplots(figsize=(6.2, 5.2), dpi=180)
    im = ax.imshow(norm, cmap="Blues", vmin=0, vmax=1)
    ax.set_title(title, fontsize=14, fontweight="bold", pad=14, color="#1c1c1c")
    ax.set_xticks(range(len(LABELS)))
    ax.set_yticks(range(len(LABELS)))
    ax.set_xticklabels(LABELS, rotation=35, ha="right", fontsize=9)
    ax.set_yticklabels(LABELS, fontsize=9)
    ax.set_xlabel("Predicted", fontsize=11, labelpad=8)
    ax.set_ylabel("Actual", fontsize=11, labelpad=8)
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            v = int(cm[i, j])
            if v == 0:
                continue
            color = "white" if norm[i, j] > 0.55 else "#1c1c1c"
            ax.text(j, i, f"{v}", ha="center", va="center",
                    fontsize=9, color=color, fontweight="bold")
    for spine in ax.spines.values():
        spine.set_visible(False)
    fig.tight_layout()
    fig.savefig(HERE / fname, dpi=180, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)


def plot_per_class_metrics(fname):
    classes = list(CNN_PR.keys())
    short = [c.replace("WALKING_", "W ") for c in classes]
    precision = [CNN_PR[c][0] for c in classes]
    recall = [CNN_PR[c][1] for c in classes]
    f1 = [CNN_PR[c][2] for c in classes]

    x = np.arange(len(classes))
    w = 0.26

    fig, ax = plt.subplots(figsize=(9.0, 4.3), dpi=180)
    ax.bar(x - w, precision, width=w, label="Precision", color="#3b82f6")
    ax.bar(x,     recall,    width=w, label="Recall",    color="#10b981")
    ax.bar(x + w, f1,        width=w, label="F1",        color="#f59e0b")
    ax.set_xticks(x)
    ax.set_xticklabels(short, rotation=15, ha="right", fontsize=10)
    ax.set_ylim(0, 1.05)
    ax.set_yticks(np.linspace(0, 1, 6))
    ax.set_yticklabels([f"{v:.0%}" for v in np.linspace(0, 1, 6)])
    ax.set_title("CNN – Per-class Precision / Recall / F1 on UCI HAR Test Set",
                 fontsize=13, fontweight="bold", pad=12, color="#1c1c1c")
    ax.legend(loc="lower right", frameon=False, fontsize=10)
    ax.grid(axis="y", alpha=0.2)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for i, (p, r, f) in enumerate(zip(precision, recall, f1)):
        for dx, v in zip([-w, 0, w], [p, r, f]):
            ax.text(i + dx, v + 0.015, f"{v:.2f}", ha="center",
                    fontsize=7.5, color="#444")
    fig.tight_layout()
    fig.savefig(HERE / fname, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ---------------------------------------------------------------- pptx helpers

def find_slide_by_title(prs, title_prefix):
    for idx, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith(title_prefix):
                return idx, s
    raise LookupError(title_prefix)


def find_placeholder(slide, name_prefix=None, idx=None, contains=None):
    for sh in slide.shapes:
        if idx is not None and getattr(sh, "placeholder_format", None) and \
                sh.placeholder_format.idx == idx:
            return sh
        if name_prefix and sh.name.startswith(name_prefix):
            return sh
        if contains and sh.has_text_frame and contains in sh.text_frame.text:
            return sh
    return None


def set_bullets(shape, bullets, sizes=None, bold_firsts=None):
    """Replace text frame content with bullets.

    bullets: list of (level:int, text:str)
    sizes:   optional list of Pt sizes aligned with bullets
    """
    tf = shape.text_frame
    tf.clear()
    if sizes is None:
        sizes = [None] * len(bullets)
    if bold_firsts is None:
        bold_firsts = [False] * len(bullets)

    for i, ((level, text), size, boldf) in enumerate(zip(bullets, sizes, bold_firsts)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        if size:
            run.font.size = Pt(size)
        if boldf:
            run.font.bold = True


def duplicate_slide(prs, src_slide):
    """Deep-copy a slide's shapes into a new slide using the same layout."""
    new_slide = prs.slides.add_slide(src_slide.slide_layout)

    # Wipe anything the layout auto-added so we don't end up with duplicates
    # on top of the copied shapes.
    sp_tree = new_slide.shapes._spTree
    for sp in list(sp_tree):
        if sp.tag.endswith("}sp") or sp.tag.endswith("}pic") or \
                sp.tag.endswith("}graphicFrame") or sp.tag.endswith("}grpSp") or \
                sp.tag.endswith("}cxnSp"):
            sp_tree.remove(sp)

    for shape in src_slide.shapes:
        el = deepcopy(shape.element)
        sp_tree.append(el)
    return new_slide


def move_slide(prs, slide, new_index):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    current = None
    for el in slides:
        # Each element is a CT_SlideIdListEntry with rId matching the slide's rels.
        rId = el.get(qn("r:id"))
        rel = prs.part.rels[rId]
        if rel.target_part is slide.part:
            current = el
            break
    if current is None:
        raise RuntimeError("could not locate slide in sldIdLst")
    sldIdLst.remove(current)
    sldIdLst.insert(new_index, current)


def clear_text_frame(shape):
    tf = shape.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def set_title(slide, text, *, color=None):
    title = slide.shapes.title
    if title is None:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title = sh
                break
    if title is None:
        return
    clear_text_frame(title)
    p = title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    if color:
        run.font.color.rgb = color


def clear_content_body(slide, keep_title=True):
    """Remove text from non-title text frames on the slide."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh is not slide.shapes.title:
            try:
                if sh.placeholder_format is not None and \
                        sh.placeholder_format.idx in (0,):  # title
                    continue
            except Exception:
                pass
            tf = sh.text_frame
            if tf.text.strip() == "":
                continue
            for p in list(tf.paragraphs[1:]):
                p._p.getparent().remove(p._p)
            p0 = tf.paragraphs[0]
            for r in list(p0.runs):
                r._r.getparent().remove(r._r)


# --------------------------------------------------------------------- build

def main():
    # Render charts first so pptx can embed them.
    plot_confusion_matrix(CNN_CM, "CNN – Confusion Matrix", "cm_cnn.png")
    plot_confusion_matrix(RF_CM, "Random Forest – Confusion Matrix", "cm_rf.png")
    plot_per_class_metrics("per_class.png")

    prs = Presentation(SRC)

    # --- Slide 7: Our System (Comparison layout: fill two empty content placeholders)
    _, s7 = find_slide_by_title(prs, "Our System")
    rf_ph = None
    cnn_ph = None
    # The Comparison layout has placeholder pairs: (header, body) x 2.
    # Indices in the original deck were: Text Placeholder 2 = "Model 1: RF",
    # Content Placeholder 3 = RF body, Text Placeholder 4 = "Model 2: 1D CNN",
    # Content Placeholder 5 = CNN body.
    for sh in s7.shapes:
        if not sh.has_text_frame:
            continue
        if sh.name == "Content Placeholder 3":
            rf_ph = sh
        elif sh.name == "Content Placeholder 5":
            cnn_ph = sh
    if rf_ph is not None:
        set_bullets(rf_ph, [
            (0, "Classical ML baseline on the 561 engineered features"),
            (0, "200 trees, bootstrap + Gini split, trained in seconds"),
            (0, "Used for comparison – strong but plateaus below the CNN"),
            (0, "Test accuracy: 85.3%"),
        ], sizes=[18, 18, 18, 18], bold_firsts=[False, False, False, True])
    if cnn_ph is not None:
        set_bullets(cnn_ph, [
            (0, "1D CNN over raw 128×9 inertial windows"),
            (0, "Conv(64)→Pool→Conv(64)→Pool→Conv(64)→GAP→Dense(100)→Softmax"),
            (0, "~120k parameters, exported to ONNX (≈170 KB)"),
            (0, "Test accuracy: 90.5%"),
        ], sizes=[18, 18, 18, 18], bold_firsts=[False, False, False, True])

    # --- Slide 8: Training Pipeline (one empty content placeholder)
    _, s8 = find_slide_by_title(prs, "Training Pipeline")
    body8 = None
    for sh in s8.shapes:
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            body8 = sh
            break
    if body8 is not None:
        set_bullets(body8, [
            (0, "Data: UCI HAR – 30 subjects, 10 299 windows, 50 Hz accel + gyro"),
            (1, "Pre-split by subject into train (7 352) / test (2 947) – respected to avoid leakage"),
            (0, "Preprocessing"),
            (1, "2.56 s sliding windows, 128 samples, 50% overlap"),
            (1, "Per-channel z-score using training mean / std (9 channels)"),
            (0, "Training"),
            (1, "Random Forest – 200 trees on 561 engineered features (baseline)"),
            (1, "1D CNN – Adam 1e-3, batch 64, early stopping on val accuracy"),
            (0, "Export"),
            (1, "Keras → ONNX via tf2onnx (opset 17) → shipped to the PWA as model.onnx"),
            (1, "preprocessing.json carries the channel order + normalization stats"),
        ])

    # --- NEW Slide: Results – numbers (insert after Training Pipeline)
    template = None
    # Use slide 10 (Limitations) as a visual template – it has the dark-title band.
    _, limits_slide = find_slide_by_title(prs, "Limitations")
    template = limits_slide
    results_slide = duplicate_slide(prs, template)
    set_title(results_slide, "Experimental Results")
    # Fill the content placeholder with headline numbers
    body = None
    for sh in results_slide.shapes:
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            body = sh
            break
    if body is not None:
        set_bullets(body, [
            (0, "1D CNN – 90.5% test accuracy on held-out subjects (proposal target ≥ 90%)"),
            (0, "Random Forest baseline – 85.3% test accuracy"),
            (0, "CNN advantage concentrated on walking-type classes"),
            (1, "WALKING         F1 0.977   (RF 0.848)"),
            (1, "WALK UPSTAIRS   F1 0.913   (RF 0.824)"),
            (1, "WALK DOWNSTAIRS F1 0.907   (RF 0.876)"),
            (0, "Both models near-perfect on LAYING (F1 ≈ 1.00)"),
            (0, "Main confusion: SITTING ↔ STANDING for both models (static + similar posture)"),
        ])
    # Add the per-class metrics bar chart
    left = Inches(5.8)
    top = Inches(4.55)
    pic = results_slide.shapes.add_picture(
        str(HERE / "per_class.png"), left, top,
        width=Inches(7.2), height=Inches(2.7)
    )
    # Move the picture behind title but above background: append() already puts it last.

    # --- NEW Slide: Results – confusion matrices side by side
    cm_slide = duplicate_slide(prs, template)
    set_title(cm_slide, "Experimental Results – Confusion Matrices")
    # Remove the content placeholder body (we only want title + images + caption).
    for sh in list(cm_slide.shapes):
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            sh.element.getparent().remove(sh.element)
    # CNN matrix (left)
    cm_slide.shapes.add_picture(
        str(HERE / "cm_cnn.png"),
        Inches(0.45), Inches(1.4),
        width=Inches(6.0), height=Inches(5.4)
    )
    # RF matrix (right)
    cm_slide.shapes.add_picture(
        str(HERE / "cm_rf.png"),
        Inches(6.85), Inches(1.4),
        width=Inches(6.0), height=Inches(5.4)
    )
    # Caption under the images
    cap = cm_slide.shapes.add_textbox(
        Inches(0.5), Inches(6.85), Inches(12.4), Inches(0.5)
    )
    cap_tf = cap.text_frame
    cap_tf.word_wrap = True
    p = cap_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Rows = true label, columns = predicted. Deeper blue = higher share. "
              "CNN collapses the stairs / walking confusion; both models split SITTING and STANDING.")
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 10: Limitations – expand content (original deck had only 1 bullet)
    body10 = None
    for sh in limits_slide.shapes:
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            body10 = sh
            break
    if body10 is not None:
        set_bullets(body10, [
            (0, "Domain shift: UCI HAR was collected with the phone clipped to the waist"),
            (1, "Our users carry the phone in a pocket or hand – axes and gravity differ"),
            (1, "Result: 90.5% on UCI test ≠ 90.5% on a real iPhone in a pocket"),
            (0, "SITTING ↔ STANDING confusion – both are static, similar gravity signature"),
            (0, "Sensor sampling variability – iOS/Android fire DeviceMotion at 50–100 Hz"),
            (1, "We resample to 50 Hz to match training, but axis conventions still differ"),
            (0, "Single-subject bias – dataset covers 30 subjects, mostly young adults"),
            (0, "Future work"),
            (1, "On-device data collection mode (CSV or Supabase upload) to retrain on real phones"),
            (1, "Orientation-invariant features (signal magnitudes, rotation to principal axis)"),
            (1, "Fine-tuning or few-shot calibration per user at first launch"),
            (1, "Extend to smartwatches for continuous wear without the pocket assumption"),
        ])

    # --- Slide 11: Conclusion – replace placeholder
    _, conc_slide = find_slide_by_title(prs, "Conclusion")
    body11 = None
    for sh in conc_slide.shapes:
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            body11 = sh
            break
    if body11 is not None:
        set_bullets(body11, [
            (0, "Delivered a working smartphone HAR system end-to-end – data → model → live web demo"),
            (0, "Trained two models on UCI HAR and evaluated on held-out subjects"),
            (1, "Random Forest baseline – 85.3% accuracy"),
            (1, "1D CNN – 90.5% accuracy, meeting the proposal target"),
            (0, "Shipped as a PWA at har-demo-cis4930.vercel.app"),
            (1, "ONNX Runtime Web runs inference in the browser – no backend, no install"),
            (1, "Supabase Realtime Broadcast lets a projector view mirror the phone in a pocket"),
            (0, "Honest limitation: transferring UCI HAR to real-world pocket use degrades accuracy"),
            (0, "Next step: collect labeled data on our own phones and retrain for the target domain"),
        ])

    # --- NEW Slide: How It Works – Tech Stack & Data Flow
    how_slide = duplicate_slide(prs, template)
    set_title(how_slide, "How It Works – Tech Stack & Data Flow")
    # Replace the content placeholder with a two-column description.
    how_body = None
    for sh in how_slide.shapes:
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            how_body = sh
            break
    if how_body is not None:
        set_bullets(how_body, [
            (0, "Data flow (end-to-end, in-browser)"),
            (1, "Phone opens har-demo-cis4930.vercel.app – a Progressive Web App loads"),
            (1, "iOS/Android DeviceMotion API streams accel + gyro at 50–100 Hz"),
            (1, "Sliding 2.56 s window (128 samples × 9 channels, matches UCI HAR)"),
            (1, "ONNX Runtime Web runs the 1D CNN on-device via WebAssembly – no server"),
            (1, "Every ~0.32 s a prediction is published to Supabase Realtime Broadcast"),
            (1, "Any other device on the same URL auto-joins as a viewer and mirrors the label"),
            (0, "Tech stack"),
            (1, "Frontend: vanilla HTML + ES modules (no framework, no build step)"),
            (1, "Inference: ONNX Runtime Web (WASM backend) – model is ~170 KB"),
            (1, "Realtime: Supabase Broadcast – ephemeral pub/sub, no tables or RLS"),
            (1, "Hosting: Vercel static deployment, auto-redeployed on every git push"),
            (1, "Training: Python · TensorFlow / Keras · tf2onnx · scikit-learn (RF baseline)"),
            (1, "Source: github.com/MateoSaettone/Har-Demo"),
        ])

    # --- NEW Slide: Live Web App – UI Walkthrough (3 screenshots)
    ui_slide = duplicate_slide(prs, template)
    set_title(ui_slide, "Live Web App – UI Walkthrough")
    # Remove the content placeholder – we use picture + caption layout instead.
    for sh in list(ui_slide.shapes):
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            sh.element.getparent().remove(sh.element)

    shot_dir = HERE / "screenshots"
    shots = [
        ("ui_idle.png",            "Idle state – awaiting a tracker"),
        ("ui_laying.png",          "Live prediction streamed to a viewer"),
        ("ui_walking_details.png", "Project details expanded on the viewer"),
    ]
    # Layout math for 13.33 × 7.5 slide: 3 images across with equal gutters.
    top = Inches(1.45)
    img_w = Inches(4.10)
    img_h = Inches(2.15)   # matches 1920×1004 aspect ratio
    gutter = Inches(0.20)
    total_w = img_w * 3 + gutter * 2
    left0 = (prs.slide_width - total_w) // 2
    cap_top = top + img_h + Inches(0.10)
    for i, (fname, caption) in enumerate(shots):
        left = left0 + (img_w + gutter) * i
        ui_slide.shapes.add_picture(
            str(shot_dir / fname), left, top, width=img_w, height=img_h
        )
        cap_box = ui_slide.shapes.add_textbox(
            left, cap_top, img_w, Inches(0.55)
        )
        tf = cap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # A descriptive paragraph under the three screenshots.
    desc = ui_slide.shapes.add_textbox(
        Inches(0.7), Inches(4.55), Inches(11.9), Inches(2.4)
    )
    dtf = desc.text_frame
    dtf.word_wrap = True
    for i, (text, size, bold) in enumerate([
        ("Single URL, two auto-detected roles.", 16, True),
        ("A device becomes a tracker the moment it taps Start tracking; every other device on the URL "
         "automatically becomes a viewer and mirrors the live label in real time — no pairing, no codes, "
         "no second app. The huge centered label is sized for readability from across a room when the "
         "phone is mirrored to a projector.", 13, False),
        ("", 6, False),
        ("On-screen chrome stays minimal: team chips, expandable Project Details with course / "
         "instructor / school / dataset / model / pipeline / repo, a Debug panel exposing the raw "
         "9-channel sensor stream and per-class probabilities, and a status pill (idle · viewing · "
         "tracking) in the top-right corner.", 13, False),
    ]):
        p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor(0x1c, 0x1c, 0x1c)

    # --- NEW Slide: Thank You / Questions (appended to the end)
    thanks_slide = duplicate_slide(prs, conc_slide)
    set_title(thanks_slide, "Thank You")
    # Replace the content placeholder with a centered message.
    for sh in thanks_slide.shapes:
        if sh.has_text_frame and sh.name.startswith("Content Placeholder"):
            tf = sh.text_frame
            tf.clear()
            p0 = tf.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            r = p0.add_run()
            r.text = "Questions?"
            r.font.size = Pt(44)
            r.font.bold = True

            for line, size, bold in [
                ("", 18, False),
                ("Tomas Pastore Godoy  ·  Mateo Saettone  ·  Ethan Thompson", 20, True),
                ("CIS4930 – Wireless and Mobile Computing", 16, False),
                ("", 10, False),
                ("Live demo:  har-demo-cis4930.vercel.app", 18, False),
                ("Source:  github.com/MateoSaettone/Har-Demo", 14, False),
            ]:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = line
                run.font.size = Pt(size)
                run.font.bold = bold
            break

    # --- Reorder: put the two Results slides right after Training Pipeline (index 7)
    # Current order after adds: [0..6, 7 TrainingPipeline, 8 FinalProduct, 9 Limitations,
    #                            10 Conclusion, 11 Results-numbers, 12 Results-CM]
    # Desired: insert 11,12 at index 8 (after Training Pipeline).
    move_slide(prs, results_slide, 8)
    move_slide(prs, cm_slide, 9)
    # After the moves above we have:
    #   0..9 previous, 10 Final Product, 11 Limits, 12 Conclusion,
    #   13 How It Works, 14 UI Walkthrough, 15 Thank You
    # Slide them in after Final Product so the order is:
    #   Final Product → How It Works → UI Walkthrough → Limitations → ...
    move_slide(prs, how_slide, 11)
    move_slide(prs, ui_slide, 12)

    prs.save(OUT)
    shutil.copy2(OUT, DOWNLOADS_OUT)
    print(f"wrote {OUT}")
    print(f"wrote {DOWNLOADS_OUT}")


if __name__ == "__main__":
    main()
